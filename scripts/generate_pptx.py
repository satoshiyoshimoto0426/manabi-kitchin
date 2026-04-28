#!/usr/bin/env python3
"""Generate PowerPoint (.pptx) training slides for ManabiOps."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brand palette
BRAND = RGBColor(0x2E, 0x7D, 0x32)
BRAND_DARK = RGBColor(0x1B, 0x5E, 0x20)
ACCENT = RGBColor(0xF5, 0x7C, 0x00)
TEXT = RGBColor(0x2C, 0x2C, 0x2C)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
BG_LIGHT = RGBColor(0xFF, 0xFD, 0xF7)
CARD_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
CARD_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
CARD_BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
CARD_RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Yu Gothic UI"  # fallback system font on Windows/Mac

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG_LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_title_bar(slide, title):
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9),
                                 Inches(0.15), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(slide, Inches(0.75), Inches(0.85), Inches(12), Inches(0.6),
             title, size=30, bold=True, color=BRAND_DARK)
    # underline
    u = slide.shapes.add_connector(1, Inches(0.5), Inches(1.45),
                                   Inches(12.8), Inches(1.45))
    u.line.color.rgb = ACCENT
    u.line.width = Pt(3)


def add_card(slide, left, top, width, height, *, bg=WHITE, border=BRAND):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.08
    card.fill.solid(); card.fill.fore_color.rgb = bg
    card.line.color.rgb = border
    card.line.width = Pt(1.5)
    return card


def add_footer(slide, page, total):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(9), Inches(0.3),
             "ManabiOps © 学び舎キッチン", size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ============ SLIDE BUILDERS ============

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # Gradient-ish: two rectangles
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = CARD_GREEN_BG; bg.line.fill.background()
    right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), 0, Inches(6.7), SH)
    right.fill.solid(); right.fill.fore_color.rgb = CARD_ORANGE_BG; right.line.fill.background()

    add_text(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.4),
             "ManabiOps", size=72, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.6),
             "こども食堂の運営を、LINEひとつで。", size=26, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
             "事務作業を 月8時間 → 30分 に", size=32, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
             "学び舎キッチン 運営自動化プラットフォーム｜研修スライド v1.0   2026年4月",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    return s


def slide_agenda():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "本日のアジェンダ（15分）")
    # Left card
    add_card(s, Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.8),
             bg=CARD_GREEN_BG, border=BRAND)
    add_text(s, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.5),
             "🎯 この研修で分かること", size=22, bold=True, color=BRAND_DARK)
    add_text(s, Inches(1.1), Inches(2.7), Inches(5.2), Inches(3.5),
             "・ManabiOps が何をしてくれるか\n・日々の使い方（LINEに送るだけ）\n・困ったときの対処法\n・安心・安全の仕組み\n・費用と役割分担",
             size=18)
    # Right card
    add_card(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8),
             bg=CARD_ORANGE_BG, border=ACCENT)
    add_text(s, Inches(7.2), Inches(2.0), Inches(5.2), Inches(0.5),
             "📋 進行", size=22, bold=True, color=BRAND_DARK)
    add_text(s, Inches(7.3), Inches(2.7), Inches(5.2), Inches(3.5),
             "1. 課題と解決（3分）\n2. 機能ツアー（5分）\n3. 使い方デモ（5分）\n4. Q&A（2分）",
             size=18)


def slide_problem():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "こんなお悩み、ありませんか？")
    items = [
        ("📒 経理", "領収書の山。月末に\nエクセル入力で 月3時間"),
        ("📝 名簿", "手書き参加者名簿を\n毎回入力で 月2時間"),
        ("📷 SNS", "写真の顔ぼかしと\n投稿準備で 月3時間"),
    ]
    for i, (t, d) in enumerate(items):
        x = Inches(0.5 + i * 4.3)
        add_card(s, x, Inches(2.0), Inches(4.0), Inches(3.0),
                 bg=CARD_RED_BG, border=RGBColor(0xC6, 0x28, 0x28))
        add_text(s, x + Inches(0.2), Inches(2.2), Inches(3.6), Inches(0.6),
                 t, size=26, bold=True, color=BRAND_DARK)
        add_text(s, x + Inches(0.2), Inches(2.9), Inches(3.6), Inches(2.0),
                 d, size=17)
    # Quote box
    add_card(s, Inches(1.5), Inches(5.4), Inches(10.3), Inches(1.2),
             bg=RGBColor(0xFF, 0xF8, 0xE1), border=ACCENT)
    add_text(s, Inches(1.8), Inches(5.6), Inches(9.8), Inches(0.9),
             "⏰ 合計 月8時間 超、本来のこども対応に使えない…",
             size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_solution():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "ManabiOps が解決します")
    add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.6),
             "LINEで送る → AIが処理 → 自動で記録",
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    items = [
        ("📱", "LINEで送るだけ", "専用アプリ不要"),
        ("🤖", "AIが自動分類", "領収書・名簿・写真"),
        ("✅", "ワンタップ承認", "間違いは防ぐ"),
    ]
    for i, (e, t, d) in enumerate(items):
        x = Inches(0.7 + i * 4.2)
        add_card(s, x, Inches(2.7), Inches(3.9), Inches(3.0),
                 bg=CARD_GREEN_BG, border=BRAND)
        add_text(s, x, Inches(2.9), Inches(3.9), Inches(0.8),
                 e, size=48, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.0), Inches(3.9), Inches(0.5),
                 t, size=22, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(4.7), Inches(3.9), Inches(0.6),
                 d, size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
             "👉 月8時間 → 月30分へ",
             size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_before_after():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "Before / After")
    rows = [
        ("作業", "Before（手作業）", "After（ManabiOps）", "削減"),
        ("領収書入力", "月3時間", "LINE送信+承認 = 5分", "-95%"),
        ("参加者名簿入力", "月2時間", "写真1枚で 自動", "-100%"),
        ("SNS投稿準備", "月3時間", "写真選んで 承認のみ", "-90%"),
        ("月次集計", "月1時間", "毎月1日 自動送信", "-100%"),
        ("合計", "月9時間", "月30分", "-94%"),
    ]
    cols = [Inches(2.5), Inches(3.0), Inches(4.5), Inches(2.0)]
    x0 = Inches(0.7); y0 = Inches(1.8); row_h = Inches(0.7)
    x = x0
    for i, (a, b, c, d) in enumerate(rows):
        header = i == 0
        total_row = i == len(rows) - 1
        bg = BRAND if header else (CARD_GREEN_BG if total_row else WHITE)
        fg = WHITE if header else (BRAND_DARK if total_row else TEXT)
        xx = x0
        for w, txt in zip(cols, [a, b, c, d]):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xx, y0 + row_h * i, w, row_h)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            tf = cell.text_frame; tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(16)
            r.font.bold = header or total_row
            r.font.color.rgb = fg
            xx += w
    add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
             "※当団体の試算。作業内容により差があります。", size=12, color=MUTED)


def slide_overall():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "全体像：3つの入口・3つの出口")
    # Left
    add_card(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.2),
             bg=CARD_BLUE_BG, border=RGBColor(0x15, 0x65, 0xC0))
    add_text(s, Inches(0.8), Inches(2.0), Inches(5.4), Inches(0.5),
             "📥 入口（スタッフがLINEに送る）", size=20, bold=True, color=BRAND_DARK)
    add_text(s, Inches(0.9), Inches(2.7), Inches(5.4), Inches(3.0),
             "🧾 領収書の写真\n📝 参加者名簿の写真\n📷 活動中の写真・動画",
             size=20)
    # Right
    add_card(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.2),
             bg=CARD_GREEN_BG, border=BRAND)
    add_text(s, Inches(7.1), Inches(2.0), Inches(5.4), Inches(0.5),
             "📤 出口（AIが自動で作る）", size=20, bold=True, color=BRAND_DARK)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.4), Inches(3.0),
             "📊 Google Sheets 会計台帳\n🔐 Firestore 参加者DB（暗号化）\n📱 Instagram 投稿（顔ぼかし済）",
             size=20)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5),
             "↓ 間に人のチェック（承認ボタン）が必ず入ります ↓",
             size=18, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
             "AIに任せきりにしない、安心設計",
             size=18, color=BRAND_DARK, align=PP_ALIGN.CENTER, bold=True)


def slide_feature(title, subtitle, steps, points, quote):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, title)
    add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
             subtitle, size=14, color=MUTED)
    # Left: steps
    add_card(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0),
             bg=CARD_GREEN_BG, border=BRAND)
    add_text(s, Inches(0.8), Inches(2.1), Inches(5.4), Inches(0.5),
             "こう動きます", size=20, bold=True, color=BRAND_DARK)
    for i, step in enumerate(steps):
        y = Inches(2.7 + i * 0.55)
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.4), Inches(0.4))
        circ.fill.solid(); circ.fill.fore_color.rgb = BRAND; circ.line.fill.background()
        tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.size = Pt(14); r.font.name = FONT
        add_text(s, Inches(1.3), y, Inches(5.0), Inches(0.5),
                 step, size=15, anchor=MSO_ANCHOR.MIDDLE)
    # Right: points
    add_card(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(3.5),
             bg=CARD_ORANGE_BG, border=ACCENT)
    add_text(s, Inches(7.1), Inches(2.1), Inches(5.4), Inches(0.5),
             "ポイント", size=20, bold=True, color=BRAND_DARK)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.5), Inches(2.8),
             "\n".join("・" + p for p in points), size=15)
    # Quote
    add_card(s, Inches(6.8), Inches(5.7), Inches(6.0), Inches(1.3),
             bg=RGBColor(0xFF, 0xF8, 0xE1), border=ACCENT)
    add_text(s, Inches(7.0), Inches(5.85), Inches(5.6), Inches(1.0),
             quote, size=15, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)


def slide_approval():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "人間中心の承認フロー")
    add_text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.7),
             "AIは「下書き」を作るだけ。最終判断は必ず人間が行います。",
             size=20, color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("①AI提案", "LINEに確認カード\n例：「コープ ¥3,240\n食材費 でOK？」", CARD_BLUE_BG, RGBColor(0x15, 0x65, 0xC0)),
        ("②人が確認", "ボタンは3つ\n✅ 承認\n✏ 修正\n❌ 却下", CARD_ORANGE_BG, ACCENT),
        ("③記録", "承認後のみ\nSheets/DB に\n書き込まれる", CARD_GREEN_BG, BRAND),
    ]
    for i, (t, d, bg, br) in enumerate(cards):
        x = Inches(0.7 + i * 4.2)
        add_card(s, x, Inches(2.6), Inches(3.9), Inches(3.3), bg=bg, border=br)
        add_text(s, x, Inches(2.8), Inches(3.9), Inches(0.6),
                 t, size=22, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(3.6), Inches(3.5), Inches(2.3),
                 d, size=16, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "24時間以内の承認を推奨。リマインドも届きます。",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)


def slide_roles():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "3つの役割")
    rows = [
        ("役割", "できること", "向いている人"),
        ("Owner (代表1名)", "全機能＋メンバー招待/削除＋削除操作", "団体代表、会計責任者"),
        ("Staff (2-5名)", "写真送信・承認・閲覧", "常勤スタッフ、会計・広報担当"),
        ("Viewer (任意)", "月次サマリーの閲覧のみ", "理事、助成元、監事"),
    ]
    cols = [Inches(3.5), Inches(5.0), Inches(4.0)]
    x0 = Inches(0.5); y0 = Inches(2.0); row_h = Inches(0.8)
    for i, row in enumerate(rows):
        header = i == 0
        xx = x0
        for w, txt in zip(cols, row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xx, y0 + row_h * i, w, row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND if header else WHITE
            cell.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            tf = cell.text_frame; tf.margin_left = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if not header else PP_ALIGN.CENTER
            r = p.add_run(); r.text = txt; r.font.name = FONT
            r.font.size = Pt(16); r.font.bold = header
            r.font.color.rgb = WHITE if header else TEXT
            xx += w
    add_card(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.9),
             bg=RGBColor(0xFF, 0xF8, 0xE1), border=ACCENT)
    add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.7),
             "💡 初期は Owner 1名 + Staff 2名 がおすすめ。慣れてから広げましょう。",
             size=16, anchor=MSO_ANCHOR.MIDDLE)


def slide_daily():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "活動日1日の流れ（例）")
    cards = [
        ("🌅 開催前 (10:00)", "・食材買い出し → レシート撮影・送信\n・会場準備 → LINEに「開始」送信", CARD_GREEN_BG, BRAND),
        ("🍚 開催中 (12-14時)", "・受付で手書き名簿記入\n・同意を得て写真撮影\n・活動名簿を撮影・送信", CARD_ORANGE_BG, ACCENT),
        ("🌇 開催後 (15:00)", "・写真・動画をLINE送信\n・AIの確認カード3-4件に承認\n・Instagram投稿完了 🎉", CARD_BLUE_BG, RGBColor(0x15, 0x65, 0xC0)),
        ("📊 月初 (翌月1日)", "・朝8時に月次サマリー到着\n・Sheetsで収支を最終確認\n・助成金報告書へコピペ", CARD_GREEN_BG, BRAND),
    ]
    for i, (t, d, bg, br) in enumerate(cards):
        x = Inches(0.5 + (i % 2) * 6.3)
        y = Inches(1.8 + (i // 2) * 2.7)
        add_card(s, x, y, Inches(6.0), Inches(2.5), bg=bg, border=br)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.5),
                 t, size=20, bold=True, color=BRAND_DARK)
        add_text(s, x + Inches(0.3), y + Inches(0.75), Inches(5.6), Inches(1.7),
                 d, size=14)


def slide_cost():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "費用の目安")
    # Left table
    add_text(s, Inches(0.7), Inches(1.8), Inches(6), Inches(0.5),
             "💰 月額コスト", size=22, bold=True, color=BRAND_DARK)
    rows = [
        ("項目", "月額"),
        ("LINE Messaging API", "¥0（200通まで無料）"),
        ("Google Cloud (AI/DB)", "¥500-¥2,000"),
        ("Instagram Graph API", "¥0"),
        ("合計", "¥500-¥3,000"),
    ]
    cols = [Inches(3.5), Inches(3.0)]
    x0 = Inches(0.7); y0 = Inches(2.4); row_h = Inches(0.6)
    for i, row in enumerate(rows):
        header = i == 0
        total_row = i == len(rows) - 1
        bg = BRAND if header else (CARD_GREEN_BG if total_row else WHITE)
        fg = WHITE if header else (ACCENT if total_row else TEXT)
        xx = x0
        for w, txt in zip(cols, row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, xx, y0 + row_h * i, w, row_h)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            cell.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            tf = cell.text_frame; tf.margin_left = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = txt; r.font.name = FONT
            r.font.size = Pt(15); r.font.bold = header or total_row
            r.font.color.rgb = fg
            xx += w
    add_text(s, Inches(0.7), Inches(5.8), Inches(6), Inches(0.4),
             "※月4回開催・参加者30名/回の想定", size=11, color=MUTED)
    # Right summary
    add_card(s, Inches(7.4), Inches(2.0), Inches(5.4), Inches(4.0),
             bg=CARD_ORANGE_BG, border=ACCENT)
    add_text(s, Inches(7.6), Inches(2.2), Inches(5.0), Inches(0.6),
             "🎯 目標との比較", size=22, bold=True, color=BRAND_DARK)
    add_text(s, Inches(7.6), Inches(3.0), Inches(5.0), Inches(0.6),
             "目標：月 ¥10,000以下", size=24, bold=True, color=TEXT)
    add_text(s, Inches(7.6), Inches(3.7), Inches(5.0), Inches(0.6),
             "実績：月 ¥500-¥3,000", size=28, bold=True, color=ACCENT)
    add_text(s, Inches(7.6), Inches(4.7), Inches(5.0), Inches(1.2),
             "助成金の「運営管理費」で\n十分まかなえる水準です。", size=16, color=BRAND_DARK)


def slide_safety():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "安心・安全への配慮")
    items = [
        ("🔐 個人情報", "・氏名ハッシュ化\n・AES-256暗号化\n・閲覧権限管理", CARD_GREEN_BG, BRAND),
        ("🙈 肖像権", "・自動顔ぼかし\n・二重検出\n・投稿前人間確認", CARD_BLUE_BG, RGBColor(0x15, 0x65, 0xC0)),
        ("📜 法令遵守", "・個人情報保護法\n・こども基本法\n・保護者同意書付", CARD_ORANGE_BG, ACCENT),
    ]
    for i, (t, d, bg, br) in enumerate(items):
        x = Inches(0.5 + i * 4.3)
        add_card(s, x, Inches(2.0), Inches(4.0), Inches(4.0), bg=bg, border=br)
        add_text(s, x, Inches(2.3), Inches(4.0), Inches(0.6),
                 t, size=22, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.4), Inches(3.3), Inches(3.6), Inches(2.5),
                 d, size=16)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "子どもたちの笑顔を、責任をもって届けます。",
             size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_phases():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "導入ステップ（目安）")
    items = [
        ("🚀 Phase 1: MVP", "・LINE Bot セットアップ\n・経理自動化 稼働開始\n・所要：4-6時間（1日）", CARD_GREEN_BG, BRAND),
        ("📈 Phase 2: フル機能", "・参加者OCR 精度調整\n・月次サマリー配信\n・所要：PoC 2週間", CARD_ORANGE_BG, ACCENT),
        ("📱 Phase 3: SNS自動化", "・IGビジネス連携審査\n・顔ぼかし・動画編集\n・所要：2-4週間", CARD_BLUE_BG, RGBColor(0x15, 0x65, 0xC0)),
        ("🎯 本番運用", "・毎週の使い方はLINEだけ\n・月次で運用チェックリスト\n・運用は月30分", CARD_GREEN_BG, BRAND),
    ]
    for i, (t, d, bg, br) in enumerate(items):
        x = Inches(0.5 + (i % 2) * 6.3)
        y = Inches(1.8 + (i // 2) * 2.7)
        add_card(s, x, y, Inches(6.0), Inches(2.5), bg=bg, border=br)
        add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.5),
                 t, size=20, bold=True, color=BRAND_DARK)
        add_text(s, x + Inches(0.3), y + Inches(0.75), Inches(5.6), Inches(1.7),
                 d, size=14)


def slide_help():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "困ったときは？")
    add_card(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.0),
             bg=CARD_BLUE_BG, border=RGBColor(0x15, 0x65, 0xC0))
    add_text(s, Inches(0.8), Inches(2.2), Inches(5.4), Inches(0.6),
             "📖 ドキュメント", size=22, bold=True, color=BRAND_DARK)
    add_text(s, Inches(0.9), Inches(2.9), Inches(5.6), Inches(3.0),
             "・機能説明書 docs/features/（10冊）\n・セットアップガイド docs/setup-guide/（15章）\n・FAQ 40問\n・印刷用PDFパック",
             size=16)
    add_card(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.0),
             bg=CARD_ORANGE_BG, border=ACCENT)
    add_text(s, Inches(7.1), Inches(2.2), Inches(5.4), Inches(0.6),
             "🆘 トラブル対応", size=22, bold=True, color=BRAND_DARK)
    add_text(s, Inches(7.2), Inches(2.9), Inches(5.6), Inches(3.0),
             "・OCR失敗 → 手入力フォームへ自動誘導\n・LINE通信エラー → 3回自動再試行\n・Instagram投稿失敗 → 管理画面通知\n・月次レポートで健全性可視化",
             size=16)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
             "💬 Owner → システム担当者の順で相談",
             size=16, color=BRAND_DARK, bold=True, align=PP_ALIGN.CENTER)


def slide_summary():
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_title_bar(s, "まとめ")
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.9),
             "📱 LINEに送るだけ", size=38, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9),
             "🤖 AIが下書きを作る", size=38, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.9),
             "✅ 最終判断は人間", size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_card(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.5),
             bg=RGBColor(0xFF, 0xF8, 0xE1), border=ACCENT)
    add_text(s, Inches(1.8), Inches(5.5), Inches(9.8), Inches(1.3),
             "「こども食堂の価値は、人と人のふれあい。\n事務はAIに、想いは人に。」",
             size=20, italic_note=False, color=BRAND_DARK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide_qa():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = CARD_GREEN_BG; bg.line.fill.background()
    right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), 0, Inches(6.7), SH)
    right.fill.solid(); right.fill.fore_color.rgb = CARD_ORANGE_BG; right.line.fill.background()
    add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
             "ご質問はありますか？", size=60, bold=True, color=BRAND_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
             "💬 ありがとうございました", size=28, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4),
             "📁 docs/features/ に詳しい資料があります", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
             "ManabiOps 研修スライド v1.0 / 学び舎キッチン / 2026年4月",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)


# Fix: add_text signature shouldn't fail on italic_note keyword
def patched_add_text(slide, left, top, width, height, text, *, size=18, bold=False,
                     color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     font=FONT, italic_note=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic_note
        r.font.color.rgb = color
    return tb

# Swap add_text for the version that accepts italic_note
globals()['add_text'] = patched_add_text


# Build all slides
slide_cover()
slide_agenda()
slide_problem()
slide_solution()
slide_before_after()
slide_overall()
slide_feature(
    "機能① 経理自動化 (FR-02/03)",
    "領収書を送るだけ → 自動で Google Sheets に記帳",
    ["領収書をスマホで撮影", "LINEで送信", "AIが店名・金額・日付を読取",
     "勘定科目を自動判定（食材費/消耗品等）", "「これで記帳？」の確認が届く",
     "✅承認 → Sheets へ自動記帳"],
    ["OCR精度が低い時は手入力フォームに自動切替",
     "重複検知（同じ領収書を2回送っても1回のみ）",
     "月末自動集計（収支・科目別）",
     "助成金報告書にそのまま使える形式"],
    "💡 領収書を撮ったらすぐ送る。溜めない運用が続くコツ。"
)
slide_feature(
    "機能② 参加者管理 (FR-04/05)",
    "手書き名簿を送るだけ → 暗号化してFirestoreに保存",
    ["当日の手書き名簿を撮影", "LINEで送信", "AIが手書き文字を認識",
     "氏名は即ハッシュ化＋暗号化", "Firestoreに安全に保存"],
    ["氏名はそのまま保存しません（SHA-256ハッシュ）",
     "復号キーは Secret Manager で厳重管理",
     "AES-256-GCM 暗号化（銀行と同等）",
     "閲覧は Owner/Staff のみ",
     "保護者同意書テンプレート付属"],
    "🛡 個人情報保護法・こども基本法に準拠"
)
slide_feature(
    "機能③ SNS投稿 (FR-06/07/08)",
    "写真・動画を送るだけ → 顔ぼかし＋キャプション自動生成で Instagram へ",
    ["活動写真・動画をLINEで送信", "AIが全員の顔を自動ぼかし",
     "動画は縦向きリール編集＋BGM", "キャプションも自動生成",
     "プレビュー → ✅承認", "Instagram自動投稿"],
    ["顔ぼかしは二重チェック（MediaPipe＋OpenCV）",
     "1人でも顔が残れば投稿ブロック",
     "ハッシュタグ自動付与（#こども食堂 等）",
     "投稿前に必ず人間の目で確認"],
    "🙈 「うっかり顔バレ」をシステムで防ぐ"
)
slide_approval()
slide_roles()
slide_daily()
slide_cost()
slide_safety()
slide_phases()
slide_help()
slide_summary()
slide_qa()

# Footer / page numbers
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    add_footer(slide, i, total)

out = Path("dist/slides/ManabiOps_研修スライド.pptx")
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"✅ Saved: {out} ({out.stat().st_size // 1024} KB, {total} slides)")
